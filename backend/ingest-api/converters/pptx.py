from __future__ import annotations

from io import BytesIO

from pptx import Presentation


def extract_pptx(file_bytes: bytes) -> tuple[str, dict]:
    presentation = Presentation(BytesIO(file_bytes))
    slides_md = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        body_parts = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if not title:
                title = text
            else:
                body_parts.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes_text = [shape.text.strip() for shape in slide.notes_slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            notes = " ".join(notes_text)
        slides_md.append(
            f"## Slide {index}: {title or f'Folie {index}'}\n"
            f"{chr(10).join(body_parts).strip()}\n"
            f"> Notes: {notes.strip()}"
        )
    return "\n\n".join(slides_md).strip(), {"slide_count": len(presentation.slides)}
